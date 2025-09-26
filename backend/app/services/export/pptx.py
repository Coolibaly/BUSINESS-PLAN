# app/export/pptx.py
from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime
import os

EXPORT_DIR = "data/exports/pptx"
os.makedirs(EXPORT_DIR, exist_ok=True)

def generate_pptx(plan_title: str, sections: list[str], filename: str) -> str:
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = f"Business Plan – {plan_title}"
    title_slide.placeholders[1].text = f"Orange Bank Afrique\n{datetime.now().strftime('%Y-%m-%d')}"

    for i, section in enumerate(sections):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Section {i+1}"
        slide.placeholders[1].text = section[:1000]  # tronqué

    path = os.path.join(EXPORT_DIR, filename)
    prs.save(path)
    return path
